"""
OPAC Document Engine  (Phase 2)
================================
Extracts clean text from every major document format.
The CPU handles file I/O; the extracted text is then passed to the NPU engine.

Supported formats
-----------------
  PDF          (.pdf)         — PyMuPDF (fitz) preferred, pypdf fallback
  Word         (.docx)        — python-docx
  PowerPoint   (.pptx)        — python-pptx
  Excel        (.xlsx / .xls) — openpyxl
  Plain text   (.txt / .md / .rst / .csv)
  HTML files   (.html / .htm)
  Web URLs     (http / https) — requests + BeautifulSoup4

All readers return a DocumentResult with:
  - text  : extracted plain text (possibly chunked by caller)
  - title : best-guess document title
  - pages : number of pages / slides / sheets (0 if unknown)
  - source: original path or URL
"""

from __future__ import annotations

import re
import unicodedata
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional
from urllib.parse import urlparse

from utils.logger import get_logger
from config.settings import DOC_MAX_CHARS, WEB_TIMEOUT_SEC, WEB_MAX_CHARS

logger = get_logger("opac.documents")


# ── Result type ───────────────────────────────────────────────────────────────

@dataclass
class DocumentResult:
    text:   str
    title:  str  = ""
    pages:  int  = 0
    source: str  = ""
    format: str  = "unknown"

    @property
    def char_count(self) -> int:
        return len(self.text)

    @property
    def word_count(self) -> int:
        return len(self.text.split())

    def __str__(self):
        return (
            f"[{self.format.upper()}] {self.title or self.source} — "
            f"{self.word_count} words, {self.pages} page(s)"
        )


# ── Dispatcher ────────────────────────────────────────────────────────────────

class DocumentLoader:
    """
    Auto-detect file type and dispatch to the correct reader.
    Usage:
        loader = DocumentLoader()
        result = loader.load("report.pdf")
        result = loader.load("https://example.com/article")
    """

    def load(self, source: str) -> DocumentResult:
        source = source.strip().strip('"').strip("'")

        # URL?
        parsed = urlparse(source)
        if parsed.scheme in ("http", "https"):
            return self._load_url(source)

        path = Path(source)
        if not path.exists():
            raise FileNotFoundError(f"File not found: {path}")

        ext = path.suffix.lower()

        readers = {
            ".pdf":  self._load_pdf,
            ".docx": self._load_docx,
            ".doc":  self._load_docx,
            ".pptx": self._load_pptx,
            ".ppt":  self._load_pptx,
            ".xlsx": self._load_xlsx,
            ".xls":  self._load_xlsx,
            ".txt":  self._load_text,
            ".md":   self._load_text,
            ".rst":  self._load_text,
            ".csv":  self._load_text,
            ".html": self._load_html_file,
            ".htm":  self._load_html_file,
        }

        reader = readers.get(ext)
        if reader is None:
            # Try plain text as last resort
            logger.warning(f"Unknown extension '{ext}', trying plain text reader.")
            return self._load_text(path)

        return reader(path)

    # ── PDF ──────────────────────────────────────────────────────────────────

    def _load_pdf(self, path: Path) -> DocumentResult:
        logger.info(f"Reading PDF: {path.name}")

        # Try PyMuPDF first (faster, better table handling)
        try:
            import fitz  # PyMuPDF
            doc   = fitz.open(str(path))
            pages = len(doc)
            title = doc.metadata.get("title", "") or path.stem

            parts = []
            for page in doc:
                parts.append(page.get_text("text"))
            doc.close()

            text = _clean("\n\n".join(parts))
            logger.info(f"PDF read via PyMuPDF: {pages} pages, {len(text)} chars")
            return DocumentResult(text=_truncate(text), title=title,
                                  pages=pages, source=str(path), format="pdf")

        except ImportError:
            logger.debug("PyMuPDF not available, falling back to pypdf")

        # Fallback: pypdf
        try:
            import pypdf
            reader = pypdf.PdfReader(str(path))
            pages  = len(reader.pages)
            title  = (reader.metadata.title or path.stem) if reader.metadata else path.stem

            parts = []
            for page in reader.pages:
                t = page.extract_text()
                if t:
                    parts.append(t)

            text = _clean("\n\n".join(parts))
            logger.info(f"PDF read via pypdf: {pages} pages, {len(text)} chars")
            return DocumentResult(text=_truncate(text), title=title,
                                  pages=pages, source=str(path), format="pdf")

        except ImportError:
            raise ImportError(
                "No PDF library found. Install one:\n"
                "  pip install PyMuPDF    (recommended)\n"
                "  pip install pypdf      (fallback)"
            )

    # ── Word (.docx) ─────────────────────────────────────────────────────────

    def _load_docx(self, path: Path) -> DocumentResult:
        logger.info(f"Reading Word document: {path.name}")
        try:
            import docx as python_docx
        except ImportError:
            raise ImportError("python-docx not installed. Run: pip install python-docx")

        doc    = python_docx.Document(str(path))
        title  = path.stem
        parts  = []

        # Core properties title
        try:
            if doc.core_properties.title:
                title = doc.core_properties.title
        except Exception:
            pass

        for para in doc.paragraphs:
            if para.text.strip():
                # Preserve heading hierarchy — guard against None style
                try:
                    style_name = para.style.name if para.style else ""
                except Exception:
                    style_name = ""
                if style_name.startswith("Heading"):
                    parts.append(f"\n## {para.text.strip()}\n")
                else:
                    parts.append(para.text.strip())

        # Tables
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                parts.append("\n[Table]\n" + "\n".join(rows) + "\n")

        text = _clean("\n\n".join(parts))
        pages = max(1, len(text) // 3000)   # rough estimate
        logger.info(f"DOCX read: ~{pages} pages, {len(text)} chars")
        return DocumentResult(text=_truncate(text), title=title,
                              pages=pages, source=str(path), format="docx")

    # ── PowerPoint (.pptx) ───────────────────────────────────────────────────

    def _load_pptx(self, path: Path) -> DocumentResult:
        logger.info(f"Reading PowerPoint: {path.name}")
        try:
            from pptx import Presentation
            from pptx.util import Pt
        except ImportError:
            raise ImportError("python-pptx not installed. Run: pip install python-pptx")

        prs    = Presentation(str(path))
        slides = len(prs.slides)
        title  = path.stem
        parts  = []

        for i, slide in enumerate(prs.slides, 1):
            slide_parts = []

            # Try to get slide title
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_title = slide.shapes.title.text.strip()
                slide_parts.append(f"[Slide {i}: {slide_title}]")
            else:
                slide_parts.append(f"[Slide {i}]")

            # Extract all text from text frames
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = " ".join(run.text for run in para.runs).strip()
                    if line and line not in slide_parts:
                        slide_parts.append(line)

            # Notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                note_text = notes_tf.text.strip() if notes_tf else ""
                if note_text:
                    slide_parts.append(f"  [Notes: {note_text}]")

            parts.append("\n".join(slide_parts))

        text = _clean("\n\n".join(parts))
        logger.info(f"PPTX read: {slides} slides, {len(text)} chars")
        return DocumentResult(text=_truncate(text), title=title,
                              pages=slides, source=str(path), format="pptx")

    # ── Excel (.xlsx) ────────────────────────────────────────────────────────

    def _load_xlsx(self, path: Path) -> DocumentResult:
        logger.info(f"Reading Excel: {path.name}")
        try:
            import openpyxl
        except ImportError:
            raise ImportError("openpyxl not installed. Run: pip install openpyxl")

        wb     = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts  = []
        sheets = 0

        for sheet_name in wb.sheetnames:
            ws     = wb[sheet_name]
            rows   = []
            header = None

            for row in ws.iter_rows(values_only=True):
                # Skip entirely empty rows
                cells = [str(c).strip() if c is not None else "" for c in row]
                if not any(cells):
                    continue
                if header is None:
                    header = cells
                    rows.append("  ".join(cells))
                else:
                    rows.append("  ".join(cells))
                if len(rows) > 200:        # cap at 200 rows per sheet
                    rows.append("  … (truncated)")
                    break

            if rows:
                parts.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
                sheets += 1

        wb.close()
        text = _clean("\n\n".join(parts))
        logger.info(f"XLSX read: {sheets} sheets, {len(text)} chars")
        return DocumentResult(text=_truncate(text), title=path.stem,
                              pages=sheets, source=str(path), format="xlsx")

    # ── Plain text ───────────────────────────────────────────────────────────

    def _load_text(self, path: Path) -> DocumentResult:
        logger.info(f"Reading text file: {path.name}")
        for enc in ("utf-8", "utf-8-sig", "latin-1", "cp1252"):
            try:
                text = path.read_text(encoding=enc)
                text = _clean(text)
                lines = len(text.splitlines())
                logger.info(f"Text file read ({enc}): {lines} lines, {len(text)} chars")
                return DocumentResult(text=_truncate(text), title=path.stem,
                                      pages=max(1, lines // 50),
                                      source=str(path), format=path.suffix.lstrip("."))
            except UnicodeDecodeError:
                continue
        raise ValueError(f"Could not decode {path} with any common encoding.")

    # ── Local HTML ───────────────────────────────────────────────────────────

    def _load_html_file(self, path: Path) -> DocumentResult:
        logger.info(f"Reading HTML file: {path.name}")
        html = path.read_text(encoding="utf-8", errors="replace")
        text, title = _parse_html(html)
        title = title or path.stem
        return DocumentResult(text=_truncate(text), title=title,
                              pages=1, source=str(path), format="html")

    # ── URL / Web page ───────────────────────────────────────────────────────

    def _load_url(self, url: str) -> DocumentResult:
        logger.info(f"Fetching URL: {url}")
        try:
            import requests
            resp = requests.get(url, timeout=WEB_TIMEOUT_SEC, headers={
                "User-Agent": "Mozilla/5.0 (OPAC local agent)"
            })
            resp.raise_for_status()
            html  = resp.text
            text, title = _parse_html(html)
            text  = _truncate(text, WEB_MAX_CHARS)
            logger.info(f"Web page fetched: {len(text)} chars, title='{title}'")
            return DocumentResult(text=text, title=title,
                                  pages=1, source=url, format="web")

        except ImportError:
            raise ImportError("requests not installed. Run: pip install requests beautifulsoup4")
        except Exception as exc:
            raise RuntimeError(f"Failed to fetch URL '{url}': {exc}") from exc


# ── HTML parsing helper ───────────────────────────────────────────────────────

def _parse_html(html: str) -> tuple[str, str]:
    """Extract readable text and title from HTML string."""
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise ImportError("beautifulsoup4 not installed. Run: pip install beautifulsoup4")

    soup  = BeautifulSoup(html, "html.parser")
    title = ""

    # Title
    t = soup.find("title")
    if t:
        title = t.get_text(strip=True)

    # Remove clutter
    for tag in soup(["script", "style", "nav", "footer", "header",
                     "aside", "form", "noscript", "iframe", "svg"]):
        tag.decompose()

    # Try to find main content area
    main = (
        soup.find("main") or
        soup.find("article") or
        soup.find(id=re.compile(r"(content|main|body|article)", re.I)) or
        soup.find(class_=re.compile(r"(content|main|body|article)", re.I)) or
        soup.body or
        soup
    )

    text = main.get_text(separator="\n", strip=True) if main else soup.get_text()
    text = _clean(text)
    return text, title


# ── Text cleaning helpers ─────────────────────────────────────────────────────

def _clean(text: str) -> str:
    """Normalise whitespace and remove non-printable characters."""
    # Unicode normalise
    text = unicodedata.normalize("NFKC", text)
    # Remove control characters except newline and tab
    text = re.sub(r"[^\S\n\t]+", " ", text)
    # Collapse 3+ blank lines to 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _truncate(text: str, limit: int = DOC_MAX_CHARS) -> str:
    if len(text) <= limit:
        return text
    logger.warning(f"Text truncated from {len(text)} to {limit} chars")
    return text[:limit] + "\n\n[… document truncated — will be chunked for summarisation …]"
