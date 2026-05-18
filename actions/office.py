"""
OPAC Office Engine  (Phase 5.5 - COM + VS Code)
==================================================
Uses Windows COM (win32com) for real-time Office control.
Changes appear INSTANTLY in Word/Excel/PowerPoint.
User can manually edit at any time.
VS Code session for code file creation and editing.

Install:  pip install pywin32 python-docx python-pptx openpyxl
VS Code:  code.exe must be in PATH
"""

from __future__ import annotations
import os, re, platform, subprocess, tempfile, time
from pathlib import Path
from typing import Optional, List, Tuple
from utils.logger import get_logger

logger     = get_logger("opac.office")
IS_WINDOWS = platform.system() == "Windows"

FOLDER_ALIASES = {
    "downloads": Path.home() / "Downloads",
    "desktop":   Path.home() / "Desktop",
    "documents": Path.home() / "Documents",
    "pictures":  Path.home() / "Pictures",
    "home":      Path.home(),
}

def _resolve_path(name, folder, ext):
    base = FOLDER_ALIASES.get(folder.lower().strip(), Path.home() / "Downloads")
    base.mkdir(parents=True, exist_ok=True)
    fn = re.sub(r'[<>:"/\\|?*]', "", name).strip() or "opac_document"
    if not fn.lower().endswith(ext): fn += ext
    return base / fn

def _open_file(path):
    try:
        if IS_WINDOWS: os.startfile(str(path))
        else: subprocess.Popen(["xdg-open", str(path)])
    except Exception as e: logger.error(f"open file: {e}")


# ======================================================================
# WORD SESSION
# ======================================================================
class WordSession:
    """COM-based Word control. Falls back to python-docx if COM unavailable."""

    def __init__(self):
        self._doc = self._word = self._docx_doc = None
        self._path = None
        self._use_com = False

    @property
    def active(self): return self._doc is not None or self._docx_doc is not None

    def new(self) -> str:
        if IS_WINDOWS:
            try:
                import win32com.client as win32
                self._word = win32.Dispatch("Word.Application")
                self._word.Visible = True
                self._doc  = self._word.Documents.Add()
                self._use_com = True
                logger.info("Word COM: new document")
                return "Word is open. Say: write about X, add heading X, add bullet list about X, save as NAME in FOLDER."
            except Exception as e:
                logger.warning(f"Word COM failed: {e}")
        return self._fallback_new()

    def open_existing(self, path: str) -> str:
        p = Path(path.strip().strip('"').strip("'"))
        if not p.exists(): return f"File not found: {path}"
        self._path = p
        if IS_WINDOWS:
            try:
                import win32com.client as win32
                self._word = win32.Dispatch("Word.Application")
                self._word.Visible = True
                self._doc  = self._word.Documents.Open(str(p.absolute()))
                self._use_com = True
                return f"Opened {p.name} in Word. Ready to edit."
            except Exception as e:
                logger.warning(f"Word COM open failed: {e}")
        from docx import Document
        self._docx_doc = Document(str(p))
        _open_file(p)
        return f"Opened {p.name}. Ready to edit."

    def _fallback_new(self) -> str:
        from docx import Document
        self._docx_doc = Document()
        for para in list(self._docx_doc.paragraphs):
            if not para.text.strip():
                para._element.getparent().remove(para._element); break
        tf = tempfile.NamedTemporaryFile(suffix=".docx", delete=False)
        self._path = Path(tf.name); tf.close()
        self._save_docx()
        _open_file(self._path)
        return "Word opened (changes visible on save)."

    def _end_range(self):
        """Return a range collapsed to document end."""
        rng = self._doc.Content
        rng.Collapse(0)  # wdCollapseEnd = 0
        return rng

    def add_heading(self, text: str, level: int = 1) -> str:
        if self._use_com:
            try:
                rng = self._end_range()
                rng.InsertParagraphAfter()
                rng.Collapse(0)
                rng.Style = self._doc.Styles(f"Heading {level}")
                rng.InsertAfter(text)
                return f"Heading {level} added: {text}"
            except Exception as e: logger.error(f"Word heading: {e}")
        if self._docx_doc:
            self._docx_doc.add_heading(text, level=level)
            self._save_docx()
        return f"Heading added: {text}"

    def add_paragraph(self, text: str) -> str:
        if self._use_com:
            try:
                rng = self._end_range()
                rng.InsertParagraphAfter()
                rng.Collapse(0)
                rng.Style = self._doc.Styles("Normal")
                rng.InsertAfter(text)
                return f"Paragraph added ({len(text.split())} words)"
            except Exception as e: logger.error(f"Word paragraph: {e}")
        if self._docx_doc:
            self._docx_doc.add_paragraph(text)
            self._save_docx()
        return f"Paragraph added ({len(text.split())} words)"

    def add_bullets(self, items: List[str]) -> str:
        if self._use_com:
            try:
                for item in items:
                    rng = self._end_range()
                    rng.InsertParagraphAfter()
                    rng.Collapse(0)
                    rng.Style = self._doc.Styles("List Bullet")
                    rng.InsertAfter(item.strip())
                return f"Added {len(items)} bullet points"
            except Exception as e: logger.error(f"Word bullets: {e}")
        if self._docx_doc:
            for item in items: self._docx_doc.add_paragraph(item.strip(), style="List Bullet")
            self._save_docx()
        return f"Added {len(items)} bullet points"

    def add_numbered(self, items: List[str]) -> str:
        if self._use_com:
            try:
                for item in items:
                    rng = self._end_range()
                    rng.InsertParagraphAfter()
                    rng.Collapse(0)
                    rng.Style = self._doc.Styles("List Number")
                    rng.InsertAfter(item.strip())
                return f"Added {len(items)} numbered items"
            except Exception as e: logger.error(f"Word numbered: {e}")
        if self._docx_doc:
            for item in items: self._docx_doc.add_paragraph(item.strip(), style="List Number")
            self._save_docx()
        return f"Added {len(items)} numbered items"

    def add_table(self, rows: int, cols: int, headers: List[str] = None) -> str:
        if self._use_com:
            try:
                rng = self._end_range()
                rng.InsertParagraphAfter()
                rng.Collapse(0)
                tbl = self._doc.Tables.Add(rng, rows, cols)
                tbl.Style = "Table Grid"
                if headers:
                    for i, h in enumerate(headers[:cols]):
                        tbl.Cell(1, i+1).Range.Text = h
                return f"Added {rows}x{cols} table"
            except Exception as e: logger.error(f"Word table: {e}")
        if self._docx_doc:
            t = self._docx_doc.add_table(rows=rows, cols=cols)
            t.style = "Table Grid"
            if headers:
                for i, h in enumerate(headers[:cols]): t.rows[0].cells[i].text = h
            self._save_docx()
        return f"Added {rows}x{cols} table"

    def add_page_break(self) -> str:
        if self._use_com:
            try:
                rng = self._end_range()
                rng.InsertBreak(7)  # wdPageBreak = 7
                return "Page break added"
            except Exception as e: logger.error(f"Word page break: {e}")
        if self._docx_doc:
            self._docx_doc.add_page_break(); self._save_docx()
        return "Page break added"

    def save(self, name: str, folder: str) -> Tuple[bool, str]:
        path = _resolve_path(name, folder, ".docx")
        try:
            if self._use_com and self._doc:
                self._doc.SaveAs2(str(path.absolute()))
                return True, f"Saved: {path.name} in {path.parent.name}"
            elif self._docx_doc:
                self._docx_doc.save(str(path)); _open_file(path)
                return True, f"Saved: {path.name} in {path.parent.name}"
        except Exception as e: return False, f"Save failed: {e}"
        return False, "Nothing to save"

    def close(self) -> str:
        try:
            if self._use_com and self._doc: self._doc.Close(SaveChanges=0)
        except Exception: pass
        self._doc = self._word = self._docx_doc = None
        return "Word session closed."

    def summary(self) -> str:
        if self._use_com and self._doc:
            try: return f"Word: {self._doc.Paragraphs.Count} paragraphs (COM live)"
            except Exception: pass
        if self._docx_doc:
            n = len([p for p in self._docx_doc.paragraphs if p.text.strip()])
            return f"Word: {n} paragraphs"
        return "Word: no document"

    def _save_docx(self):
        if self._docx_doc and self._path:
            try: self._docx_doc.save(str(self._path))
            except Exception: pass


# ======================================================================
# POWERPOINT SESSION
# ======================================================================
class PowerPointSession:
    def __init__(self):
        self._prs = self._pptx = self._app = None
        self._path = None; self._use_com = False; self._count = 0

    @property
    def active(self): return self._prs is not None or self._pptx is not None

    def new(self, title: str = "") -> str:
        if IS_WINDOWS:
            try:
                import win32com.client as win32
                self._app = win32.Dispatch("PowerPoint.Application")
                self._app.Visible = True
                self._prs = self._app.Presentations.Add()
                self._use_com = True
                if title: self._com_add_title_slide(title)
                return f"PowerPoint is open: '{title}'. Say: create next slide about X."
            except Exception as e: logger.warning(f"PPTX COM failed: {e}")
        return self._fallback_new(title)

    def open_existing(self, path: str) -> str:
        p = Path(path.strip().strip('"').strip("'"))
        if not p.exists(): return f"File not found: {path}"
        self._path = p
        if IS_WINDOWS:
            try:
                import win32com.client as win32
                self._app = win32.Dispatch("PowerPoint.Application")
                self._app.Visible = True
                self._prs  = self._app.Presentations.Open(str(p.absolute()))
                self._use_com = True
                self._count = self._prs.Slides.Count
                return f"Opened {p.name} ({self._count} slides). Ready."
            except Exception as e: logger.warning(f"PPTX COM open failed: {e}")
        from pptx import Presentation
        self._pptx = Presentation(str(p))
        self._count = len(self._pptx.slides); _open_file(p)
        return f"Opened {p.name} ({self._count} slides)."

    def _fallback_new(self, title: str = "") -> str:
        from pptx import Presentation
        self._pptx = Presentation()
        if title:
            s = self._pptx.slides.add_slide(self._pptx.slide_layouts[0])
            s.shapes.title.text = title; self._count = 1
        tf = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        self._path = Path(tf.name); tf.close()
        self._pptx.save(str(self._path)); _open_file(self._path)
        return f"PowerPoint opened: '{title}'."

    def _com_add_title_slide(self, title: str, subtitle: str = ""):
        slide = self._prs.Slides.Add(1, 1)  # ppLayoutTitle=1
        slide.Shapes(1).TextFrame.TextRange.Text = title
        if subtitle:
            try: slide.Shapes(2).TextFrame.TextRange.Text = subtitle
            except Exception: pass
        self._count += 1

    def _com_add_slide(self, layout: int = 2) -> object:
        return self._prs.Slides.Add(self._prs.Slides.Count + 1, layout)

    def add_content_slide(self, title: str, content: str) -> str:
        self._count += 1
        if self._use_com:
            try:
                s = self._com_add_slide(2)
                s.Shapes(1).TextFrame.TextRange.Text = title
                s.Shapes(2).TextFrame.TextRange.Text = content
                return f"Slide {self._count}: {title}"
            except Exception as e: logger.error(f"PPTX content slide: {e}")
        if self._pptx:
            s = self._pptx.slides.add_slide(self._pptx.slide_layouts[1])
            s.shapes.title.text = title; s.placeholders[1].text = content
            self._save_pptx()
        return f"Slide {self._count}: {title}"

    def add_bullet_slide(self, title: str, bullets: List[str]) -> str:
        self._count += 1
        if self._use_com:
            try:
                s = self._com_add_slide(2)
                s.Shapes(1).TextFrame.TextRange.Text = title
                tf = s.Shapes(2).TextFrame
                tf.TextRange.Text = bullets[0] if bullets else ""
                for b in bullets[1:]: tf.TextRange.InsertAfter("\r" + b)
                return f"Bullet slide {self._count}: {title}"
            except Exception as e: logger.error(f"PPTX bullet slide: {e}")
        if self._pptx:
            s = self._pptx.slides.add_slide(self._pptx.slide_layouts[1])
            s.shapes.title.text = title
            tf = s.placeholders[1].text_frame
            tf.text = bullets[0] if bullets else ""
            for b in bullets[1:]:
                p = tf.add_paragraph(); p.text = b
            self._save_pptx()
        return f"Bullet slide {self._count}: {title}"

    def add_section_slide(self, title: str) -> str:
        self._count += 1
        if self._use_com:
            try:
                s = self._com_add_slide(11)  # ppLayoutTitleOnly=11
                s.Shapes(1).TextFrame.TextRange.Text = title
                return f"Section slide {self._count}: {title}"
            except Exception as e: logger.error(f"PPTX section: {e}")
        if self._pptx:
            try: layout = self._pptx.slide_layouts[2]
            except: layout = self._pptx.slide_layouts[0]
            s = self._pptx.slides.add_slide(layout)
            s.shapes.title.text = title; self._save_pptx()
        return f"Section slide {self._count}: {title}"

    def add_blank_slide(self) -> str:
        self._count += 1
        if self._use_com:
            try: self._com_add_slide(12); return f"Blank slide {self._count}"
            except Exception as e: logger.error(f"PPTX blank: {e}")
        if self._pptx:
            try: layout = self._pptx.slide_layouts[6]
            except: layout = self._pptx.slide_layouts[1]
            self._pptx.slides.add_slide(layout); self._save_pptx()
        return f"Blank slide {self._count}"

    def save(self, name: str, folder: str) -> Tuple[bool, str]:
        path = _resolve_path(name, folder, ".pptx")
        try:
            if self._use_com and self._prs:
                self._prs.SaveAs(str(path.absolute()))
                return True, f"Saved: {path.name} in {path.parent.name}"
            elif self._pptx:
                self._pptx.save(str(path)); _open_file(path)
                return True, f"Saved: {path.name} in {path.parent.name}"
        except Exception as e: return False, f"Save failed: {e}"
        return False, "Nothing to save"

    def close(self) -> str:
        try:
            if self._use_com and self._prs: self._prs.Close()
        except Exception: pass
        self._prs = self._pptx = self._app = None
        return "PowerPoint session closed."

    def summary(self) -> str: return f"PowerPoint: {self._count} slides"

    def _save_pptx(self):
        if self._pptx and self._path:
            try: self._pptx.save(str(self._path))
            except Exception: pass


# ======================================================================
# EXCEL SESSION
# ======================================================================
class ExcelSession:
    def __init__(self):
        self._wb = self._ws = self._xl = None
        self._oxl = self._ows = None
        self._path = None; self._use_com = False; self._row = 1

    @property
    def active(self): return self._wb is not None or self._oxl is not None

    def new(self) -> str:
        if IS_WINDOWS:
            try:
                import win32com.client as win32
                self._xl = win32.Dispatch("Excel.Application")
                self._xl.Visible = True
                self._wb = self._xl.Workbooks.Add()
                self._ws = self._wb.ActiveSheet
                self._use_com = True
                return "Excel is open. Say: add headers X Y Z, add row X Y Z, add total row, save as NAME in FOLDER."
            except Exception as e: logger.warning(f"Excel COM failed: {e}")
        return self._fallback_new()

    def open_existing(self, path: str) -> str:
        p = Path(path.strip().strip('"').strip("'"))
        if not p.exists(): return f"File not found: {path}"
        self._path = p
        if IS_WINDOWS:
            try:
                import win32com.client as win32
                self._xl = win32.Dispatch("Excel.Application")
                self._xl.Visible = True
                self._wb = self._xl.Workbooks.Open(str(p.absolute()))
                self._ws = self._wb.ActiveSheet
                self._use_com = True
                self._row = self._ws.UsedRange.Rows.Count + 1
                return f"Opened {p.name}. Next row: {self._row}."
            except Exception as e: logger.warning(f"Excel COM open failed: {e}")
        from openpyxl import load_workbook
        self._oxl = load_workbook(str(p))
        self._ows = self._oxl.active
        self._row = self._ows.max_row + 1; _open_file(p)
        return f"Opened {p.name}. Next row: {self._row}."

    def _fallback_new(self) -> str:
        from openpyxl import Workbook
        self._oxl = Workbook(); self._ows = self._oxl.active
        self._ows.title = "Sheet1"; self._row = 1
        tf = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False)
        self._path = Path(tf.name); tf.close()
        self._save_oxl(); _open_file(self._path)
        return "Excel opened (changes visible on save)."

    def add_headers(self, headers: List[str]) -> str:
        if self._use_com:
            try:
                for c, h in enumerate(headers, 1):
                    cell = self._ws.Cells(self._row, c)
                    cell.Value = h; cell.Font.Bold = True
                    cell.Interior.Color = 0xC44244
                    cell.Font.Color = 0xFFFFFF
                self._row += 1
                return f"Headers: {', '.join(headers)}"
            except Exception as e: logger.error(f"Excel headers: {e}")
        if self._oxl:
            from openpyxl.styles import Font, PatternFill, Alignment
            for c, h in enumerate(headers, 1):
                cell = self._ows.cell(row=self._row, column=c, value=h)
                cell.font      = Font(bold=True, color="FFFFFF")
                cell.fill      = PatternFill("solid", fgColor="4472C4")
                cell.alignment = Alignment(horizontal="center")
            self._row += 1; self._save_oxl()
        return f"Headers: {', '.join(headers)}"

    def add_row(self, values: list) -> str:
        if self._use_com:
            try:
                for c, v in enumerate(values, 1):
                    self._ws.Cells(self._row, c).Value = v
                self._row += 1
                return f"Row {self._row-1} added"
            except Exception as e: logger.error(f"Excel row: {e}")
        if self._oxl:
            for c, v in enumerate(values, 1):
                self._ows.cell(row=self._row, column=c, value=v)
            self._row += 1; self._save_oxl()
        return f"Row {self._row-1} added"

    def add_total_row(self) -> str:
        if self._row < 3: return "Not enough data"
        if self._use_com:
            try:
                cols = self._ws.UsedRange.Columns.Count
                self._ws.Cells(self._row, 1).Value = "TOTAL"
                self._ws.Cells(self._row, 1).Font.Bold = True
                for c in range(2, cols + 1):
                    addr = self._ws.Cells(1, c).Address(False, False)
                    col_letter = re.sub(r'\d', '', addr)
                    self._ws.Cells(self._row, c).Formula = f"=SUM({col_letter}2:{col_letter}{self._row-1})"
                    self._ws.Cells(self._row, c).Font.Bold = True
                self._row += 1
                return "Total row added"
            except Exception as e: logger.error(f"Excel total: {e}")
        if self._oxl:
            from openpyxl.styles import Font
            cols = self._ows.max_column
            self._ows.cell(row=self._row, column=1, value="TOTAL").font = Font(bold=True)
            for c in range(2, cols+1):
                cl = self._ows.cell(row=1, column=c).column_letter
                cell = self._ows.cell(row=self._row, column=c, value=f"=SUM({cl}2:{cl}{self._row-1})")
                cell.font = Font(bold=True)
            self._row += 1; self._save_oxl()
        return "Total row added"

    def create_sheet(self, name: str) -> str:
        if self._use_com:
            try:
                self._ws = self._wb.Sheets.Add()
                self._ws.Name = name; self._row = 1
                return f"Sheet '{name}' created"
            except Exception as e: logger.error(f"Excel sheet: {e}")
        if self._oxl:
            ws = self._oxl.create_sheet(title=name)
            self._ows = ws; self._row = 1; self._save_oxl()
        return f"Sheet '{name}' created"

    def save(self, name: str, folder: str) -> Tuple[bool, str]:
        path = _resolve_path(name, folder, ".xlsx")
        try:
            if self._use_com and self._wb:
                self._wb.SaveAs(str(path.absolute()))
                return True, f"Saved: {path.name} in {path.parent.name}"
            elif self._oxl:
                self._save_oxl(path); _open_file(path)
                return True, f"Saved: {path.name} in {path.parent.name}"
        except Exception as e: return False, f"Save failed: {e}"
        return False, "Nothing to save"

    def close(self) -> str:
        try:
            if self._use_com and self._wb: self._wb.Close(SaveChanges=False)
        except Exception: pass
        self._wb = self._ws = self._xl = self._oxl = self._ows = None
        return "Excel session closed."

    def summary(self) -> str: return f"Excel: {self._row-1} rows"

    def _save_oxl(self, path=None):
        p = path or self._path
        if self._oxl and p:
            try:
                for col in self._ows.columns:
                    w = max((len(str(c.value or "")) for c in col), default=10)
                    self._ows.column_dimensions[col[0].column_letter].width = min(w+4, 40)
                self._oxl.save(str(p))
            except Exception: pass


# ======================================================================
# VS CODE SESSION
# ======================================================================
class VSCodeSession:
    EXT_MAP = {
        "python": ".py", "javascript": ".js", "typescript": ".ts",
        "html": ".html", "css": ".css", "java": ".java",
        "c": ".c", "c++": ".cpp", "cpp": ".cpp", "csharp": ".cs", "c#": ".cs",
        "go": ".go", "rust": ".rs", "ruby": ".rb", "php": ".php",
        "sql": ".sql", "bash": ".sh", "powershell": ".ps1",
        "markdown": ".md", "json": ".json", "yaml": ".yml",
        "xml": ".xml", "text": ".txt",
    }
    HEADERS = {
        "python":     "# {name}.py\n# Created by OPAC\n\n",
        "javascript": "// {name}.js\n// Created by OPAC\n\n",
        "typescript": "// {name}.ts\n// Created by OPAC\n\n",
        "html": "<!DOCTYPE html>\n<html lang=\"en\">\n<head>\n  <meta charset=\"UTF-8\">\n  <title>{name}</title>\n</head>\n<body>\n\n</body>\n</html>\n",
        "css":  "/* {name}.css - Created by OPAC */\n\n",
        "java": "// {name}.java\npublic class {classname} {{\n\n}}\n",
    }
    EDITOR_CMDS = {
        "vscode": "code", "vs code": "code", "code": "code",
        "cursor": "cursor", "sublime": "subl", "notepad++": "notepad++",
    }

    def __init__(self, editor: str = "vscode"):
        self._cmd       = self.EDITOR_CMDS.get(editor.lower(), "code")
        self._file_path: Optional[Path] = None
        self._language  = "python"

    @property
    def active(self): return self._file_path is not None

    def open_editor(self, path: str = "") -> str:
        import shutil
        cmd = shutil.which(self._cmd) or self._cmd
        try:
            if path:
                p = Path(path.strip().strip('"').strip("'"))
                flags = subprocess.DETACHED_PROCESS if IS_WINDOWS else 0
                subprocess.Popen([cmd, str(p)], creationflags=flags)
                if p.is_file():
                    self._file_path = p
                    self._language  = self._detect_lang(p.suffix)
                return f"Opened {p.name} in {self._cmd}"
            else:
                flags = subprocess.DETACHED_PROCESS if IS_WINDOWS else 0
                subprocess.Popen([cmd], creationflags=flags)
                return f"{self._cmd} launched. Say: create python file called X."
        except Exception as e:
            return f"Cannot open {self._cmd}: {e}. Is it installed and in PATH?"

    def create_file(self, name: str, language: str, folder: str = "documents") -> str:
        ext  = self.EXT_MAP.get(language.lower(), ".txt")
        path = _resolve_path(name, folder, ext)
        self._language = language.lower(); self._file_path = path
        header = self.HEADERS.get(self._language, "# {name}\n# Created by OPAC\n\n")
        header = header.format(name=path.stem, classname=path.stem.title())
        path.write_text(header, encoding="utf-8")
        import shutil
        cmd = shutil.which(self._cmd) or self._cmd
        try:
            flags = subprocess.DETACHED_PROCESS if IS_WINDOWS else 0
            subprocess.Popen([cmd, str(path)], creationflags=flags)
        except Exception as e: logger.error(f"VS Code create: {e}")
        return f"Created {path.name} and opened in {self._cmd}. Say: write a function that does X."

    def open_existing(self, path: str) -> str:
        p = Path(path.strip().strip('"').strip("'"))
        if not p.exists(): return f"File not found: {path}"
        self._file_path = p; self._language = self._detect_lang(p.suffix)
        import shutil
        cmd = shutil.which(self._cmd) or self._cmd
        try:
            flags = subprocess.DETACHED_PROCESS if IS_WINDOWS else 0
            subprocess.Popen([cmd, str(p)], creationflags=flags)
        except Exception as e: logger.error(f"VS Code open: {e}")
        return f"Opened {p.name} in {self._cmd}"

    def append_code(self, code: str) -> str:
        if not self._file_path: return "No file open. Say 'create python file called X' first."
        try:
            current = self._file_path.read_text(encoding="utf-8")
            sep     = "\n\n" if current.strip() else ""
            self._file_path.write_text(current.rstrip("\n") + sep + code + "\n", encoding="utf-8")
            lines = code.count("\n") + 1
            return f"Added {lines} lines to {self._file_path.name}"
        except Exception as e: return f"Write error: {e}"

    def replace_content(self, code: str) -> str:
        if not self._file_path: return "No file open."
        try:
            self._file_path.write_text(code + "\n", encoding="utf-8")
            return f"File {self._file_path.name} updated"
        except Exception as e: return f"Write error: {e}"

    def read_content(self) -> str:
        if not self._file_path or not self._file_path.exists(): return ""
        return self._file_path.read_text(encoding="utf-8")

    def run_file(self) -> str:
        if not self._file_path: return "No file to run."
        runners = {
            ".py": ["python", str(self._file_path)],
            ".js": ["node",   str(self._file_path)],
            ".rb": ["ruby",   str(self._file_path)],
        }
        cmd = runners.get(self._file_path.suffix.lower())
        if not cmd: return f"Cannot run {self._file_path.suffix} files directly."
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
            out = r.stdout[:500] or r.stderr[:500] or "(no output)"
            return f"Output:\n{out}"
        except subprocess.TimeoutExpired: return "Timed out."
        except Exception as e: return f"Run error: {e}"

    def close(self) -> str:
        self._file_path = None; return "VS Code session closed."

    def summary(self) -> str:
        if self._file_path:
            try: lines = self._file_path.read_text().count("\n")
            except: lines = 0
            return f"VS Code: {self._file_path.name} ({lines} lines)"
        return "VS Code: no file"

    def _detect_lang(self, ext: str) -> str:
        return {v: k for k, v in self.EXT_MAP.items()}.get(ext.lower(), "text")


# ======================================================================
# BROWSER SESSION
# ======================================================================
class BrowserSession:
    """
    Browser control via subprocess — no Playwright, no asyncio conflicts.

    Strategy:
      - open()        : launch browser normally via subprocess
      - navigate()    : pass URL as argument to browser exe
      - search_*()    : build URL and pass to browser exe
      - new_tab()     : open browser with --new-tab flag
      - whatsapp_send(): open WhatsApp Web URL, then pyautogui

    Why not Playwright:
      httpx (used by Wikipedia API) leaves an asyncio event loop running.
      Playwright sync API cannot start inside an existing loop.
      subprocess avoids ALL of these issues and works with real profiles.
    """

    EXE_PATHS = {
        "chrome": [
            r"C:\Program Files\Google\Chrome\Application\chrome.exe",
            r"C:\Program Files (x86)\Google\Chrome\Application\chrome.exe",
        ],
        "edge": [
            r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe",
            r"C:\Program Files\Microsoft\Edge\Application\msedge.exe",
        ],
        "brave": [
            r"C:\Program Files\BraveSoftware\Brave-Browser\Application\brave.exe",
            r"C:\Program Files (x86)\BraveSoftware\Brave-Browser\Application\brave.exe",
        ],
    }
    USER_DATA_DIRS = {
        "chrome": Path.home() / "AppData/Local/Google/Chrome/User Data",
        "edge":   Path.home() / "AppData/Local/Microsoft/Edge/User Data",
        "brave":  Path.home() / "AppData/Local/BraveSoftware/Brave-Browser/User Data",
    }

    def __init__(self):
        self._exe  = ""
        self._name = ""
        self._page = None   # Playwright page (only for summarize tab)
        self._pw   = None
        self._context = None
        self._browser = None

    @property
    def active(self):
        # Active as long as we have an exe to use
        return bool(self._exe)

    def _find_exe(self, name: str) -> str:
        for path in self.EXE_PATHS.get(name, []):
            if Path(path).exists():
                return path
        import shutil
        return shutil.which(name) or shutil.which(name + ".exe") or ""

    def _popen(self, args: list):
        """Launch process detached so it doesn't block OPAC."""
        flags = subprocess.DETACHED_PROCESS | subprocess.CREATE_NEW_PROCESS_GROUP if IS_WINDOWS else 0
        subprocess.Popen(args, creationflags=flags,
                         stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)

    def _open_url(self, url: str):
        """Open a URL in the current browser."""
        if self._exe:
            self._popen([self._exe, url])
        else:
            import webbrowser
            webbrowser.open(url)

    # ── Open browser ──────────────────────────────────────────────────────────

    def open(self, name: str = "chrome", profile: str = "") -> Tuple[bool, str]:
        """
        Launch browser normally via subprocess with optional profile.
        No Playwright involved — user's real browser with real profile.
        """
        name = name.lower().strip()
        exe  = self._find_exe(name)
        if not exe:
            return False, f"{name} not found. Is it installed?"

        self._exe  = exe
        self._name = name

        args = [exe]
        if profile:
            pdir = self._resolve_profile_dir(name, profile)
            udir = self.USER_DATA_DIRS.get(name)
            if pdir and udir and udir.exists():
                args += [f"--user-data-dir={udir}", f"--profile-directory={pdir}"]

        self._popen(args)
        profile_msg = f" (profile: {profile})" if profile else ""
        return True, f"Opened {name.title()}{profile_msg}"

    def open_normal(self, name: str = "chrome", profile: str = "") -> Tuple[bool, str]:
        """Alias for open() — same behaviour."""
        return self.open(name, profile)

    # ── Navigation ────────────────────────────────────────────────────────────

    def navigate(self, url: str) -> Tuple[bool, str]:
        if not url.startswith("http"):
            url = "https://" + url
        self._open_url(url)
        return True, f"Opening {url}"

    def search_google(self, q: str) -> Tuple[bool, str]:
        import urllib.parse
        url = "https://www.google.com/search?q=" + urllib.parse.quote_plus(q)
        self._open_url(url)
        return True, f"Searching Google: {q}"

    def search_youtube(self, q: str) -> Tuple[bool, str]:
        import urllib.parse
        url = "https://www.youtube.com/results?search_query=" + urllib.parse.quote_plus(q)
        self._open_url(url)
        return True, f"Searching YouTube: {q}"

    def new_tab(self, url: str = "") -> Tuple[bool, str]:
        if url and url.startswith("http"):
            self._open_url(url)
            return True, f"Opened tab: {url}"
        if self._exe:
            self._popen([self._exe, "--new-tab"])
        return True, "New tab opened"

    def get_page_text(self) -> str:
        """Try to get page text via CDP/Playwright (optional, for summarize tab)."""
        if self._page:
            try:
                return self._page.evaluate("""() => {
                    ['script','style','nav','footer','header']
                        .forEach(t=>document.querySelectorAll(t).forEach(e=>e.remove()));
                    return document.body.innerText||'';
                }""")[:3000]
            except Exception:
                pass
        return ""

    # ── WhatsApp Web ──────────────────────────────────────────────────────────

    def whatsapp_send(self, contact: str, message: str) -> Tuple[bool, str]:
        """
        Open WhatsApp Web and send message via pyautogui.
        Opens wa.me deeplink which goes directly to the chat if contact is a number,
        otherwise opens the search.
        """
        # Try wa.me deeplink for phone numbers
        if re.match(r"^\+?[\d\s\-]{7,}$", contact):
            phone = re.sub(r"[\s\-]", "", contact).lstrip("+")
            import urllib.parse
            url = f"https://wa.me/{phone}?text={urllib.parse.quote_plus(message)}"
            self._open_url(url)
            return True, f"WhatsApp deeplink opened for {contact}"

        # For contact names: open WhatsApp Web then use pyautogui
        self._open_url("https://web.whatsapp.com")
        time.sleep(6)

        try:
            import pyautogui, pyperclip
            pyautogui.FAILSAFE = False
            # Search for contact using keyboard shortcut
            pyautogui.hotkey("ctrl", "alt", "n") if IS_WINDOWS else None
            time.sleep(1)
            # Click search
            pyautogui.hotkey("ctrl", "f")
            time.sleep(1)
            pyperclip.copy(contact)
            pyautogui.hotkey("ctrl", "v")
            time.sleep(2)
            pyautogui.press("enter")
            time.sleep(1.5)
            # Click message box and type
            w, h = pyautogui.size()
            pyautogui.click(w // 2, int(h * 0.92))
            time.sleep(0.8)
            pyperclip.copy(message)
            pyautogui.hotkey("ctrl", "v")
            time.sleep(0.3)
            pyautogui.press("enter")
            return True, f"WhatsApp message sent to {contact}"
        except ImportError:
            return True, f"WhatsApp Web opened (install pyautogui for auto-send)"
        except Exception as e:
            return False, f"WhatsApp error: {e}"

    # ── Close ─────────────────────────────────────────────────────────────────

    def close(self) -> str:
        self._exe = self._name = ""
        if self._browser:
            try: self._browser.close()
            except Exception: pass
        if self._pw:
            try: self._pw.stop()
            except Exception: pass
        self._browser = self._context = self._page = self._pw = None
        return "Browser session closed"

    # ── CDP connect (for summarize tab only) ──────────────────────────────────

    def connect_cdp(self, port: int = 9222) -> Tuple[bool, str]:
        """Connect to already-running browser via CDP for page reading."""
        for p in [port, 9223, 9224]:
            try:
                from playwright.sync_api import sync_playwright
                pw = sync_playwright().start()
                browser = pw.chromium.connect_over_cdp(f"http://localhost:{p}")
                self._pw = pw; self._browser = browser
                ctx = browser.contexts[0] if browser.contexts else browser.new_context()
                self._context = ctx
                pages = ctx.pages
                self._page = pages[0] if pages else ctx.new_page()
                return True, f"Connected on port {p}"
            except Exception:
                try: pw.stop()
                except Exception: pass
                continue
        return False, "No browser found on debug ports 9222-9224"

    # ── Profile resolution ────────────────────────────────────────────────────

    def _resolve_profile_dir(self, browser: str, name: str) -> Optional[str]:
        udir = self.USER_DATA_DIRS.get(browser)
        if not udir or not udir.exists():
            return None
        nl = name.lower().strip()
        if nl in ("default", "1", "first", "main", "personal"):
            return "Default"
        m = re.search(r"\d+", name)
        if m:
            n = int(m.group())
            return "Default" if n == 1 else f"Profile {n - 1}"
        try:
            import json
            for d in udir.iterdir():
                if d.is_dir() and (d.name == "Default" or d.name.startswith("Profile")):
                    prefs = d / "Preferences"
                    if prefs.exists():
                        pn = json.loads(prefs.read_text(errors="ignore")
                             ).get("profile", {}).get("name", "").lower()
                        if nl in pn or pn in nl:
                            return d.name
        except Exception:
            pass
        return None

    def _find_profile(self, browser, name):
        return self._resolve_profile_dir(browser, name)


class MessagingEngine:
    """
    Send messages via Viber, WhatsApp desktop, Telegram, Discord.
    Uses EnumWindows for robust window finding (handles partial titles).
    Waits up to 15 seconds for app to load after opening.
    """

    @staticmethod
    def _find_window(app_name: str, timeout: float = 15.0) -> int:
        """
        Find a window containing app_name in its title.
        Waits up to timeout seconds for it to appear.
        Returns HWND or 0 if not found.
        """
        if not IS_WINDOWS:
            return 0
        import ctypes

        found = [0]
        app_lower = app_name.lower()

        WNDENUMPROC = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.c_int, ctypes.c_int)

        def enum_callback(hwnd, lparam):
            length = ctypes.windll.user32.GetWindowTextLengthW(hwnd)
            if length > 0:
                buf = ctypes.create_unicode_buffer(length + 1)
                ctypes.windll.user32.GetWindowTextW(hwnd, buf, length + 1)
                title = buf.value.lower()
                if app_lower in title and ctypes.windll.user32.IsWindowVisible(hwnd):
                    found[0] = hwnd
                    return False   # stop enumeration
            return True

        deadline = time.time() + timeout
        while time.time() < deadline and found[0] == 0:
            ctypes.windll.user32.EnumWindows(WNDENUMPROC(enum_callback), 0)
            if found[0] == 0:
                time.sleep(0.5)

        return found[0]

    def send_desktop(self, app: str, contact: str, message: str) -> Tuple[bool, str]:
        """Send message via desktop app using pyautogui."""
        try:
            import pyautogui
            try:
                import pyperclip; has_clip = True
            except ImportError:
                has_clip = False
        except ImportError:
            raise ImportError("pip install pyautogui pyperclip")

        pyautogui.FAILSAFE = False

        print(f"  [OPAC] Looking for {app} window ...", flush=True)
        hwnd = self._find_window(app, timeout=15.0)

        if not hwnd:
            return False, (
                f"{app} window not found after 15 seconds. "
                f"Make sure {app} is open and visible (not minimized to system tray)."
            )

        # Focus the window
        import ctypes
        u32 = ctypes.windll.user32
        u32.ShowWindow(hwnd, 9)         # SW_RESTORE
        u32.SetForegroundWindow(hwnd)
        time.sleep(1.0)

        app_lower = app.lower()

        # Open new chat / search
        if app_lower == "viber":
            # Viber: Ctrl+N opens new message dialog
            pyautogui.hotkey("ctrl", "n")
            time.sleep(1.5)
        elif app_lower == "telegram":
            pyautogui.hotkey("ctrl", "k")
            time.sleep(1.0)
        elif app_lower in ("whatsapp",):
            pyautogui.hotkey("ctrl", "f")
            time.sleep(1.0)
        else:
            pyautogui.hotkey("ctrl", "f")
            time.sleep(1.0)

        # Type contact name
        pyautogui.hotkey("ctrl", "a")
        time.sleep(0.2)
        if has_clip:
            pyperclip.copy(contact)
            pyautogui.hotkey("ctrl", "v")
        else:
            for ch in contact:
                pyautogui.press(ch)
                time.sleep(0.05)
        time.sleep(1.5)

        # Press Enter to open the chat
        pyautogui.press("enter")
        time.sleep(1.0)

        # Click the message input area (bottom centre of window)
        # Get window rect for precise clicking
        try:
            import ctypes
            rect = ctypes.wintypes.RECT()
            ctypes.windll.user32.GetWindowRect(hwnd, ctypes.byref(rect))
            win_x = (rect.left + rect.right) // 2
            win_y = rect.top + int((rect.bottom - rect.top) * 0.92)
            pyautogui.click(win_x, win_y)
        except Exception:
            w, h = pyautogui.size()
            pyautogui.click(w // 2, int(h * 0.92))
        time.sleep(0.8)

        # Type message
        if has_clip:
            pyperclip.copy(message)
            pyautogui.hotkey("ctrl", "v")
        else:
            for ch in message:
                pyautogui.press(ch)
                time.sleep(0.04)
        time.sleep(0.4)
        pyautogui.press("enter")

        return True, f"Message sent to {contact} via {app}"


# ======================================================================
# SIMPLE TEXT EDITOR SESSION (Notepad, CMD, PowerShell)
# ======================================================================
class SimpleEditorSession:
    """
    Voice control for Notepad, CMD, and PowerShell.

    Notepad — same functionality as Word:
      write about X in N words   → NPU generates, writes to file
      add heading X              → adds formatted heading
      add bullet list about X    → NPU generates bullets
      append X                   → NPU generates, appends
      save as NAME in FOLDER     → saves .txt file

    CMD / PowerShell — commands run IN the open terminal window:
      run command X             → types X into terminal, presses Enter
      Result appears in the terminal, not in OPAC.
    """

    def __init__(self, app: str = "notepad"):
        self._app      = app.lower().strip()
        self._path: Optional[Path] = None
        self._content  = ""
        self._hwnd     = 0

    @property
    def active(self): return True

    def new(self) -> str:
        if self._app == "notepad":
            import tempfile
            tf = tempfile.NamedTemporaryFile(
                suffix=".txt", delete=False, mode="w", encoding="utf-8")
            self._path    = Path(tf.name)
            self._content = ""
            tf.close()
            _open_file(self._path)
            time.sleep(1.0)
            return ("Notepad is open. Commands: "
                    "write about X, add heading X, add bullet list about X, "
                    "append X, save as NAME in FOLDER.")
        elif self._app == "cmd":
            if IS_WINDOWS:
                subprocess.Popen("cmd.exe",
                    creationflags=subprocess.CREATE_NEW_CONSOLE)
            else:
                subprocess.Popen(["x-terminal-emulator"])
            time.sleep(1.5)
            return "CMD opened. Say: run command X — it will type and run in the terminal."
        elif self._app in ("powershell", "ps"):
            if IS_WINDOWS:
                subprocess.Popen("powershell.exe",
                    creationflags=subprocess.CREATE_NEW_CONSOLE)
            else:
                subprocess.Popen(["x-terminal-emulator", "-e", "pwsh"])
            time.sleep(1.5)
            return "PowerShell opened. Say: run command X — it runs in the terminal."
        return f"Unknown app: {self._app}"

    # ── Notepad content methods ────────────────────────────────────────────────

    def write(self, text: str) -> str:
        """Overwrite file content."""
        if self._app != "notepad": return "Write only for Notepad."
        self._content = text
        self._sync()
        return f"Written ({len(text.split())} words)"

    def append_text(self, text: str) -> str:
        """Append content with blank line separator."""
        if self._app != "notepad": return "Append only for Notepad."
        sep = "\n\n" if self._content.strip() else ""
        self._content = self._content.rstrip("\n") + sep + text
        self._sync()
        return f"Appended ({len(text.split())} words)"

    def add_heading(self, text: str) -> str:
        """Add a formatted heading (underlined with ===)."""
        if self._app != "notepad": return "Headings only for Notepad."
        bar = "=" * max(len(text), 30)
        heading = f"\n{bar}\n{text.upper()}\n{bar}\n"
        sep = "\n" if self._content.strip() else ""
        self._content = self._content.rstrip("\n") + sep + heading
        self._sync()
        return f"Heading added: {text}"

    def add_bullets(self, items: List[str]) -> str:
        """Add bullet list."""
        if self._app != "notepad": return "Bullets only for Notepad."
        bullet_text = "\n".join(f"  • {item.strip()}" for item in items if item.strip())
        sep = "\n\n" if self._content.strip() else ""
        self._content = self._content.rstrip("\n") + sep + bullet_text
        self._sync()
        return f"Added {len(items)} bullet points"

    # ── CMD / PowerShell — run IN the open terminal ────────────────────────────

    def run_in_terminal(self, command: str) -> str:
        """
        Type and run a command in the open CMD or PowerShell window.
        Result appears in the terminal itself, not in OPAC.
        """
        if not IS_WINDOWS:
            return self._run_subprocess(command)

        # Find terminal window (with retry up to 10s)
        terminal_keywords = {
            "cmd":        ["cmd.exe", "command prompt"],
            "powershell": ["windows powershell", "powershell"],
            "ps":         ["windows powershell", "powershell"],
        }
        keywords = terminal_keywords.get(self._app, [self._app])

        hwnd = 0
        import ctypes
        WNDENUMPROC = ctypes.WINFUNCTYPE(ctypes.c_bool, ctypes.c_int, ctypes.c_int)
        found = [0]

        for _ in range(20):   # up to 10 seconds
            found[0] = 0
            def enum_cb(hwnd, _):
                length = ctypes.windll.user32.GetWindowTextLengthW(hwnd)
                if length > 0:
                    buf = ctypes.create_unicode_buffer(length + 1)
                    ctypes.windll.user32.GetWindowTextW(hwnd, buf, length + 1)
                    title = buf.value.lower()
                    if (any(k in title for k in keywords)
                            and ctypes.windll.user32.IsWindowVisible(hwnd)):
                        found[0] = hwnd
                        return False
                return True
            ctypes.windll.user32.EnumWindows(WNDENUMPROC(enum_cb), 0)
            hwnd = found[0]
            if hwnd: break
            time.sleep(0.5)

        if not hwnd:
            return f"Terminal window not found. Showing output here instead:\n{self._run_subprocess(command)}"

        # Focus window
        u32 = ctypes.windll.user32
        u32.ShowWindow(hwnd, 9)
        u32.SetForegroundWindow(hwnd)
        time.sleep(0.5)

        # Paste command using clipboard (handles special characters)
        try:
            import pyperclip
            pyperclip.copy(command)
            import pyautogui
            pyautogui.hotkey("ctrl", "v")
        except Exception:
            try:
                import pyautogui
                pyautogui.typewrite(command, interval=0.04)
            except Exception:
                return self._run_subprocess(command)

        import pyautogui
        pyautogui.press("enter")
        return f"Command sent to {self._app} terminal: {command}"

    def _run_subprocess(self, command: str) -> str:
        """Fallback: run in subprocess and return output to OPAC."""
        if self._app == "cmd":
            runner = ["cmd", "/c", command]
        else:
            runner = ["powershell", "-NoProfile", "-Command", command]
        try:
            r = subprocess.run(runner, capture_output=True, text=True, timeout=30)
            return (r.stdout + r.stderr).strip()[:800] or "(no output)"
        except Exception as e:
            return f"Error: {e}"

    # ── Save ──────────────────────────────────────────────────────────────────

    def save(self, name: str, folder: str) -> Tuple[bool, str]:
        if self._app != "notepad":
            return False, f"Save not applicable for {self._app}"
        path = _resolve_path(name, folder, ".txt")
        try:
            path.write_text(self._content, encoding="utf-8")
            self._path = path
            _open_file(path)
            return True, f"Saved: {path.name} in {path.parent.name}"
        except Exception as e:
            return False, f"Save failed: {e}"

    def close(self) -> str:
        self._content = ""
        return f"{self._app.title()} session closed."

    def summary(self) -> str:
        if self._app == "notepad" and self._content:
            return f"Notepad: {len(self._content.split())} words"
        return f"{self._app.title()}: active"

    def _sync(self):
        """Write content to temp file so Notepad refreshes."""
        if self._path:
            try:
                self._path.write_text(self._content, encoding="utf-8")
            except Exception:
                pass




# ======================================================================
# OFFICE MANAGER
# ======================================================================
class OfficeManager:
    def __init__(self):
        self.word:    Optional[WordSession]         = None
        self.pptx:    Optional[PowerPointSession]   = None
        self.excel:   Optional[ExcelSession]        = None
        self.browser: Optional[BrowserSession]      = None
        self.vscode:  Optional[VSCodeSession]       = None
        self.editor:  Optional[SimpleEditorSession] = None
        self.msg      = MessagingEngine()

    # ── Word ──────────────────────────────────────────────────────────────────

    def start_word(self, path: str = "") -> str:
        self.word = WordSession()
        return self.word.open_existing(path) if path else self.word.new()

    def word_cmd(self, cmd: str, content: str, extra: dict = None) -> str:
        if not self.word: return "No Word document. Say 'open word'."
        extra = extra or {}; w = self.word
        if cmd == "heading":    return w.add_heading(content, extra.get("level",1))
        if cmd == "subheading": return w.add_heading(content, level=2)
        if cmd == "paragraph":  return w.add_paragraph(content)
        if cmd == "bullets":
            items = [l.strip() for l in content.split("\n") if l.strip()]
            return w.add_bullets(items)
        if cmd == "numbered":
            items = [l.strip() for l in content.split("\n") if l.strip()]
            return w.add_numbered(items)
        if cmd == "table":
            return w.add_table(extra.get("rows",3), extra.get("cols",3), extra.get("headers",[]))
        if cmd == "page_break": return w.add_page_break()
        if cmd == "save":
            ok,msg = w.save(extra.get("name","document"), extra.get("folder","downloads"))
            if ok: self.word = None
            return msg
        if cmd == "close":
            msg = w.close(); self.word = None; return msg
        if cmd == "status": return w.summary()
        return f"Unknown: {cmd}"

    # ── PowerPoint ────────────────────────────────────────────────────────────

    def start_pptx(self, title: str = "", path: str = "") -> str:
        self.pptx = PowerPointSession()
        return self.pptx.open_existing(path) if path else self.pptx.new(title)

    def pptx_cmd(self, cmd: str, content: str, extra: dict = None) -> str:
        if not self.pptx: return "No presentation. Say 'open powerpoint'."
        extra = extra or {}; p = self.pptx
        if cmd == "content_slide": return p.add_content_slide(extra.get("title","Slide"), content)
        if cmd == "bullet_slide":
            items = [l.strip() for l in content.split("\n") if l.strip()]
            return p.add_bullet_slide(extra.get("title","Slide"), items)
        if cmd == "title_slide":   return p.add_title_slide(content, extra.get("subtitle",""))
        if cmd == "section_slide": return p.add_section_slide(content)
        if cmd == "blank_slide":   return p.add_blank_slide()
        if cmd == "save":
            ok,msg = p.save(extra.get("name","presentation"), extra.get("folder","downloads"))
            if ok: self.pptx = None
            return msg
        if cmd == "close":
            msg = p.close(); self.pptx = None; return msg
        if cmd == "status": return p.summary()
        return f"Unknown: {cmd}"

    # ── Excel ─────────────────────────────────────────────────────────────────

    def start_excel(self, path: str = "") -> str:
        self.excel = ExcelSession()
        return self.excel.open_existing(path) if path else self.excel.new()

    def excel_cmd(self, cmd: str, content: str, extra: dict = None) -> str:
        if not self.excel: return "No Excel workbook. Say 'open excel'."
        extra = extra or {}; e = self.excel
        if cmd == "headers":
            h = [x.strip() for x in re.split(r"[,\s]+", content) if x.strip()]
            return e.add_headers(h)
        if cmd == "row":    return e.add_row(_parse_row(content))
        if cmd == "total":  return e.add_total_row()
        if cmd == "sheet":  return e.create_sheet(content)
        if cmd == "save":
            ok,msg = e.save(extra.get("name","spreadsheet"), extra.get("folder","downloads"))
            if ok: self.excel = None
            return msg
        if cmd == "close":
            msg = e.close(); self.excel = None; return msg
        if cmd == "status": return e.summary()
        return f"Unknown: {cmd}"

    # ── VS Code ───────────────────────────────────────────────────────────────

    def start_vscode(self, editor: str = "vscode", path: str = "") -> str:
        if not self.vscode:
            self.vscode = VSCodeSession(editor)
        return self.vscode.open_existing(path) if path else self.vscode.open_editor()

    def vscode_cmd(self, cmd: str, content: str, extra: dict = None) -> str:
        if not self.vscode: self.vscode = VSCodeSession()
        extra = extra or {}; vs = self.vscode
        if cmd == "create":
            return vs.create_file(extra.get("name","untitled"),
                                  extra.get("language","python"),
                                  extra.get("folder","documents"))
        if cmd == "open":    return vs.open_existing(content)
        if cmd == "append":  return vs.append_code(content)
        if cmd == "replace": return vs.replace_content(content)
        if cmd == "run":     return vs.run_file()
        if cmd == "read":    return vs.read_content()
        if cmd == "close":
            msg = vs.close(); self.vscode = None; return msg
        if cmd == "status":  return vs.summary()
        return f"Unknown: {cmd}"

    # ── Simple Editor (Notepad / CMD / PowerShell) ────────────────────────────

    def start_editor(self, app: str = "notepad") -> str:
        self.editor = SimpleEditorSession(app)
        return self.editor.new()

    def editor_cmd(self, cmd: str, content: str, extra: dict = None) -> str:
        if not self.editor: return "No editor open. Say 'open notepad' first."
        extra = extra or {}
        if cmd == "write":    return self.editor.write(content)
        if cmd == "append":   return self.editor.append_text(content)
        if cmd == "heading":  return self.editor.add_heading(content)
        if cmd == "bullets":
            items = [l.strip() for l in content.split("\n") if l.strip()]
            return self.editor.add_bullets(items)
        if cmd == "run":      return self.editor.run_in_terminal(content)
        if cmd == "save":
            ok, msg = self.editor.save(extra.get("name","notes"), extra.get("folder","downloads"))
            if ok: self.editor = None
            return msg
        if cmd == "close":
            msg = self.editor.close(); self.editor = None; return msg
        if cmd == "status":   return self.editor.summary()
        return f"Unknown editor command: {cmd}"

    # ── Browser ───────────────────────────────────────────────────────────────

    def start_browser(self, name: str = "chrome", profile: str = "",
                      for_automation: bool = True) -> Tuple[bool, str]:
        if for_automation:
            if self.browser and self.browser.active:
                try: self.browser.close()
                except Exception: pass
            self.browser = BrowserSession()
            return self.browser.open(name, profile)
        else:
            b = BrowserSession()
            return b.open_normal(name, profile)

    def browser_cmd(self, cmd: str, content: str = "") -> Tuple[bool, str]:
        if not self.browser or not self.browser.active:
            return False, "No browser open."
        b = self.browser
        if cmd == "navigate":       return b.navigate(content)
        if cmd == "search":         return b.search_google(content)
        if cmd == "search_youtube": return b.search_youtube(content)
        if cmd == "new_tab":        return b.new_tab(content)
        if cmd == "read":
            t = b.get_page_text(); return True, t if t else "Page empty"
        if cmd == "whatsapp":
            parts = content.split("|",1)
            return b.whatsapp_send(parts[0].strip(), parts[1].strip() if len(parts)>1 else "")
        if cmd == "close":
            msg = b.close(); self.browser = None; return True, msg
        return False, f"Unknown: {cmd}"

    # ── Messaging ─────────────────────────────────────────────────────────────

    def send_message(self, app: str, contact: str, message: str) -> Tuple[bool, str]:
        if app.lower() == "whatsapp" and self.browser and self.browser.active:
            return self.browser.whatsapp_send(contact, message)
        return self.msg.send_desktop(app, contact, message)

    # ── Status / close all ────────────────────────────────────────────────────

    def status(self) -> str:
        parts = []
        if self.word:                            parts.append(self.word.summary())
        if self.pptx:                            parts.append(self.pptx.summary())
        if self.excel:                           parts.append(self.excel.summary())
        if self.vscode:                          parts.append(self.vscode.summary())
        if self.editor:                          parts.append(self.editor.summary())
        if self.browser and self.browser.active: parts.append("Browser: open (OPAC controlled)")
        return "\n  ".join(parts) if parts else "No active sessions."

    def close_all(self):
        for s in [self.word, self.pptx, self.excel, self.vscode, self.editor]:
            if s:
                try: s.close()
                except Exception: pass
        self.word = self.pptx = self.excel = self.vscode = self.editor = None
        if self.browser:
            try: self.browser.close()
            except Exception: pass
        self.browser = None


# ── Helpers ───────────────────────────────────────────────────────────────────

def _parse_row(text: str) -> list:
    import re as _re
    parts = _re.split(r"[,\t]+|\s{2,}", text.strip())
    result = []
    for p in parts:
        p = p.strip()
        if not p: continue
        try:    result.append(int(p))
        except ValueError:
            try:    result.append(float(p))
            except: result.append(p)
    return result


def parse_save(text: str) -> Tuple[str, str]:
    import re as _re
    m = _re.search(r"save\s+(?:as\s+)?(.+?)\s+(?:in|to|into|at)\s+(\w+)\s*$", text, _re.I)
    if m: return m.group(1).strip(), m.group(2).lower()
    m = _re.search(r"save\s+(?:as\s+)?(.+)$", text, _re.I)
    if m: return m.group(1).strip(), "downloads"
    return "document", "downloads"


def parse_slide(text: str) -> Tuple[str, str]:
    import re as _re
    m = _re.search(r"title\s+(.+?)\s+and\s+content\s+(.+)$", text, _re.I)
    if m: return m.group(1).strip(), m.group(2).strip()
    m = _re.search(r"(?:about|on|covering|for)\s+(.+)$", text, _re.I)
    if m: t = m.group(1).strip(); return t.title(), t
    m = _re.search(r"called\s+(.+)$", text, _re.I)
    if m: t = m.group(1).strip(); return t, t
    return "Slide", text