"""
OPAC Test Suite — Phase 1 & 2
==============================
Run from the project root:
    python tests/test_phase1_2.py

Tests are split into two groups:
  - Phase 1: NPU engine (requires OpenVINO + NPU driver)
  - Phase 2: Document engine (runs on any machine with the libraries)

Phase 2 tests run without the NPU by using a MockEngine that returns
a fixed string.  This lets you verify document reading independently.
"""

from __future__ import annotations

import os
import sys
import textwrap
import tempfile
import unittest
from pathlib import Path

# Make sure project root is on the path
ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT))

from documents.loader import DocumentLoader, DocumentResult
from utils.chunker import chunk_text


# ══════════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════════

class MockEngine:
    """Stands in for NPUEngine so document tests run without the NPU."""

    loaded = True
    device = "MOCK"

    def load(self): pass
    def unload(self): pass

    def generate(self, prompt: str, max_new_tokens=512, streamer_callback=None) -> str:
        snippet = prompt[:80].replace("\n", " ")
        return f"[MockEngine summary of: {snippet}…]"

    def build_summarize_prompt(self, content): return f"SUMMARIZE:{content}"
    def build_chunk_prompt(self, content):     return f"CHUNK:{content}"
    def build_combine_prompt(self, summaries): return f"COMBINE:{summaries}"
    def build_prompt(self, msg, history=""):   return f"CHAT:{msg}"


def _tmp_file(suffix: str, content: bytes) -> Path:
    f = tempfile.NamedTemporaryFile(suffix=suffix, delete=False)
    f.write(content)
    f.close()
    return Path(f.name)


# ══════════════════════════════════════════════════════════════════════════════
# Phase 2 — Document Engine tests  (no NPU required)
# ══════════════════════════════════════════════════════════════════════════════

class TestChunker(unittest.TestCase):
    """Chunker splits and preserves content correctly."""

    def test_short_text_no_chunk(self):
        text   = "Hello world. " * 10
        chunks = chunk_text(text, max_chars=500)
        self.assertEqual(len(chunks), 1)
        self.assertIn("Hello world", chunks[0])

    def test_long_text_multiple_chunks(self):
        text   = ("A sentence that is roughly twenty chars. " * 100)
        chunks = chunk_text(text, max_chars=200, overlap=20)
        self.assertGreater(len(chunks), 1)
        for c in chunks:
            self.assertLessEqual(len(c), 220)   # allow small overage at boundary

    def test_chunk_content_coverage(self):
        """All original words should appear in at least one chunk."""
        words  = [f"word{i}" for i in range(500)]
        text   = " ".join(words)
        chunks = chunk_text(text, max_chars=300, overlap=30)
        all_text = " ".join(chunks)
        for w in words:
            self.assertIn(w, all_text)

    def test_empty_text(self):
        chunks = chunk_text("", max_chars=100)
        self.assertEqual(chunks, [])

    def test_exact_boundary(self):
        text   = "x" * 1000
        chunks = chunk_text(text, max_chars=1000, overlap=0)
        self.assertEqual(len(chunks), 1)


class TestDocumentLoaderText(unittest.TestCase):
    """Plain-text file reading."""

    def setUp(self):
        self.loader = DocumentLoader()

    def test_read_txt(self):
        content  = "This is a test document.\nIt has multiple lines.\nLine three."
        tmp      = _tmp_file(".txt", content.encode("utf-8"))
        try:
            result = self.loader.load(str(tmp))
            self.assertIsInstance(result, DocumentResult)
            self.assertIn("test document", result.text)
            self.assertEqual(result.format, "txt")
        finally:
            tmp.unlink()

    def test_read_md(self):
        content = "# Heading\n\nSome **markdown** content.\n\n- item 1\n- item 2"
        tmp     = _tmp_file(".md", content.encode("utf-8"))
        try:
            result = self.loader.load(str(tmp))
            self.assertIn("Heading", result.text)
            self.assertEqual(result.format, "md")
        finally:
            tmp.unlink()

    def test_utf8_content(self):
        content = "Namaste नमस्ते. Special: café, résumé, naïve."
        tmp     = _tmp_file(".txt", content.encode("utf-8"))
        try:
            result = self.loader.load(str(tmp))
            self.assertIn("Namaste", result.text)
        finally:
            tmp.unlink()

    def test_empty_file(self):
        tmp = _tmp_file(".txt", b"")
        try:
            result = self.loader.load(str(tmp))
            self.assertEqual(result.text.strip(), "")
        finally:
            tmp.unlink()

    def test_file_not_found(self):
        with self.assertRaises(FileNotFoundError):
            self.loader.load("/nonexistent/path/file.pdf")


class TestDocumentLoaderDocx(unittest.TestCase):
    """Word document reading."""

    def setUp(self):
        self.loader = DocumentLoader()

    def test_read_docx(self):
        try:
            import docx as python_docx
        except ImportError:
            self.skipTest("python-docx not installed")

        doc  = python_docx.Document()
        doc.add_heading("Test Report", level=1)
        doc.add_paragraph("This is the introduction paragraph.")
        doc.add_heading("Section 2", level=2)
        doc.add_paragraph("Second section content goes here.")

        tmp = _tmp_file(".docx", b"")
        doc.save(str(tmp))
        try:
            result = self.loader.load(str(tmp))
            self.assertIn("Test Report", result.text)
            self.assertIn("introduction paragraph", result.text)
            self.assertEqual(result.format, "docx")
        finally:
            tmp.unlink()

    def test_docx_with_table(self):
        try:
            import docx as python_docx
        except ImportError:
            self.skipTest("python-docx not installed")

        doc   = python_docx.Document()
        table = doc.add_table(rows=2, cols=3)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Score"
        table.cell(0, 2).text = "Grade"
        table.cell(1, 0).text = "Alice"
        table.cell(1, 1).text = "95"
        table.cell(1, 2).text = "A"

        tmp = _tmp_file(".docx", b"")
        doc.save(str(tmp))
        try:
            result = self.loader.load(str(tmp))
            self.assertIn("Alice", result.text)
            self.assertIn("95", result.text)
        finally:
            tmp.unlink()


class TestDocumentLoaderPptx(unittest.TestCase):
    """PowerPoint reading."""

    def setUp(self):
        self.loader = DocumentLoader()

    def test_read_pptx(self):
        try:
            from pptx import Presentation
            from pptx.util import Inches
        except ImportError:
            self.skipTest("python-pptx not installed")

        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "My Presentation"
        slide.placeholders[1].text = "Key point one\nKey point two"

        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Slide Two"
        slide2.placeholders[1].text = "Another important point."

        tmp = _tmp_file(".pptx", b"")
        prs.save(str(tmp))
        try:
            result = self.loader.load(str(tmp))
            self.assertIn("My Presentation", result.text)
            self.assertIn("Key point", result.text)
            self.assertEqual(result.pages, 2)
            self.assertEqual(result.format, "pptx")
        finally:
            tmp.unlink()

    def test_pptx_notes(self):
        try:
            from pptx import Presentation
        except ImportError:
            self.skipTest("python-pptx not installed")

        prs   = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide with Notes"
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = "Speaker note: remember to mention the budget."

        tmp = _tmp_file(".pptx", b"")
        prs.save(str(tmp))
        try:
            result = self.loader.load(str(tmp))
            self.assertIn("budget", result.text)
        finally:
            tmp.unlink()


class TestDocumentLoaderXlsx(unittest.TestCase):
    """Excel reading."""

    def setUp(self):
        self.loader = DocumentLoader()

    def test_read_xlsx(self):
        try:
            import openpyxl
        except ImportError:
            self.skipTest("openpyxl not installed")

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Sales"
        ws.append(["Month", "Revenue", "Units"])
        ws.append(["January", 50000, 120])
        ws.append(["February", 62000, 145])

        tmp = _tmp_file(".xlsx", b"")
        wb.save(str(tmp))
        try:
            result = self.loader.load(str(tmp))
            self.assertIn("Revenue", result.text)
            self.assertIn("January", result.text)
            self.assertIn("62000", result.text)
            self.assertEqual(result.format, "xlsx")
        finally:
            tmp.unlink()

    def test_xlsx_multiple_sheets(self):
        try:
            import openpyxl
        except ImportError:
            self.skipTest("openpyxl not installed")

        wb = openpyxl.Workbook()
        ws1 = wb.active
        ws1.title = "Sheet1"
        ws1.append(["Alpha", "Beta"])

        ws2 = wb.create_sheet("Sheet2")
        ws2.append(["Gamma", "Delta"])

        tmp = _tmp_file(".xlsx", b"")
        wb.save(str(tmp))
        try:
            result = self.loader.load(str(tmp))
            self.assertIn("Alpha", result.text)
            self.assertIn("Gamma", result.text)
        finally:
            tmp.unlink()


class TestDocumentLoaderHTML(unittest.TestCase):
    """HTML file and web content parsing."""

    def setUp(self):
        self.loader = DocumentLoader()

    def test_read_html_file(self):
        html = b"""
        <html><head><title>Test Page</title></head>
        <body>
          <nav>Navigation stuff</nav>
          <main>
            <h1>Main Heading</h1>
            <p>This is the important content of the page.</p>
            <p>Second paragraph with more information.</p>
          </main>
          <footer>Footer content</footer>
        </body></html>
        """
        tmp = _tmp_file(".html", html)
        try:
            result = self.loader.load(str(tmp))
            self.assertIn("Main Heading", result.text)
            self.assertIn("important content", result.text)
            self.assertEqual(result.title, "Test Page")
            # Nav and footer should be stripped
            self.assertNotIn("Navigation stuff", result.text)
        finally:
            tmp.unlink()

    def test_html_strips_scripts(self):
        html = b"""
        <html><body>
          <script>alert('bad')</script>
          <style>.x { color: red }</style>
          <p>Real content here.</p>
        </body></html>
        """
        tmp = _tmp_file(".html", html)
        try:
            result = self.loader.load(str(tmp))
            self.assertNotIn("alert", result.text)
            self.assertNotIn("color: red", result.text)
            self.assertIn("Real content", result.text)
        finally:
            tmp.unlink()


class TestSummarizer(unittest.TestCase):
    """Summarizer correctly routes to MockEngine."""

    def setUp(self):
        from core.summarizer import Summarizer
        self.engine     = MockEngine()
        self.summarizer = Summarizer(self.engine)

    def test_short_text(self):
        result = self.summarizer.summarize_text("Short text.", stream=False)
        self.assertIsInstance(result, str)
        self.assertGreater(len(result), 0)

    def test_long_text_chunks(self):
        long_text = "This is a sentence with ten words exactly here. " * 300
        result    = self.summarizer.summarize_text(long_text, stream=False)
        self.assertIsInstance(result, str)

    def test_empty_text(self):
        result = self.summarizer.summarize_text("", stream=False)
        self.assertIn("empty", result.lower())

    def test_summarize_txt_file(self):
        content = "The quarterly results exceeded all projections significantly. " * 20
        tmp     = _tmp_file(".txt", content.encode())
        try:
            result = self.summarizer.summarize_file(str(tmp), stream=False)
            self.assertIsInstance(result, str)
        finally:
            tmp.unlink()


# ══════════════════════════════════════════════════════════════════════════════
# Phase 1 — NPU Engine tests  (skipped if OpenVINO not installed)
# ══════════════════════════════════════════════════════════════════════════════

class TestNPUEngineImport(unittest.TestCase):
    """NPUEngine class loads without errors even if OpenVINO is absent."""

    def test_import(self):
        from core.npu_engine import NPUEngine
        engine = NPUEngine()
        self.assertFalse(engine.loaded)

    def test_generate_without_load_raises(self):
        from core.npu_engine import NPUEngine
        engine = NPUEngine()
        with self.assertRaises(RuntimeError):
            engine.generate("hello")

    def test_prompt_building(self):
        from core.npu_engine import NPUEngine
        engine = NPUEngine()
        prompt = engine.build_summarize_prompt("Some document content.")
        self.assertIn("Some document content", prompt)
        self.assertIn("<|user|>", prompt)
        self.assertIn("<|assistant|>", prompt)

    def test_chunk_prompt(self):
        from core.npu_engine import NPUEngine
        engine = NPUEngine()
        prompt = engine.build_chunk_prompt("A section of text.")
        self.assertIn("A section of text", prompt)

    def test_combine_prompt(self):
        from core.npu_engine import NPUEngine
        engine = NPUEngine()
        prompt = engine.build_combine_prompt(["Summary one.", "Summary two."])
        self.assertIn("Summary one", prompt)
        self.assertIn("Summary two", prompt)


def _openvino_available() -> bool:
    try:
        import openvino       # noqa
        import openvino_genai # noqa
        return True
    except ImportError:
        return False


@unittest.skipUnless(
    _openvino_available(),
    "Skipping NPU tests — OpenVINO not installed"
)
class TestNPUEngineWithOpenVINO(unittest.TestCase):
    """Tests that run only when OpenVINO is installed (on your Acer Swift AI)."""

    def test_npu_visible(self):
        import openvino as ov
        core    = ov.Core()
        devices = core.available_devices
        self.assertIn("NPU", devices, f"NPU not found. Available: {devices}")

    def test_model_dir_exists_after_setup(self):
        from config.settings import DEFAULT_MODEL_DIR
        if not DEFAULT_MODEL_DIR.exists():
            self.skipTest("Model not downloaded yet — run 'python opac.py --setup' first")
        self.assertTrue(any(DEFAULT_MODEL_DIR.iterdir()))

    def test_load_and_generate(self):
        from config.settings import DEFAULT_MODEL_DIR, INFERENCE_DEVICE
        if not DEFAULT_MODEL_DIR.exists():
            self.skipTest("Model not downloaded yet")

        from core.npu_engine import NPUEngine
        engine = NPUEngine(model_dir=DEFAULT_MODEL_DIR, device=INFERENCE_DEVICE)
        try:
            engine.load()
            self.assertTrue(engine.loaded)
            result = engine.generate(
                engine.build_prompt("Reply with exactly: OPAC_TEST_OK"),
                max_new_tokens=20,
            )
            self.assertIsInstance(result, str)
            self.assertGreater(len(result), 0)
            print(f"\n      NPU response: {result}")
        finally:
            engine.unload()


# ══════════════════════════════════════════════════════════════════════════════
# Runner
# ══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    print("=" * 60)
    print(" OPAC Test Suite — Phase 1 & 2")
    print("=" * 60)
    loader = unittest.TestLoader()
    suite  = unittest.TestSuite()

    # Always run Phase 2 (document) tests
    suite.addTests(loader.loadTestsFromTestCase(TestChunker))
    suite.addTests(loader.loadTestsFromTestCase(TestDocumentLoaderText))
    suite.addTests(loader.loadTestsFromTestCase(TestDocumentLoaderDocx))
    suite.addTests(loader.loadTestsFromTestCase(TestDocumentLoaderPptx))
    suite.addTests(loader.loadTestsFromTestCase(TestDocumentLoaderXlsx))
    suite.addTests(loader.loadTestsFromTestCase(TestDocumentLoaderHTML))
    suite.addTests(loader.loadTestsFromTestCase(TestSummarizer))

    # Phase 1 — NPU (skips gracefully without OpenVINO)
    suite.addTests(loader.loadTestsFromTestCase(TestNPUEngineImport))
    suite.addTests(loader.loadTestsFromTestCase(TestNPUEngineWithOpenVINO))

    runner = unittest.TextTestRunner(verbosity=2)
    result = runner.run(suite)
    sys.exit(0 if result.wasSuccessful() else 1)
